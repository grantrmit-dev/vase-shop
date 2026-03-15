#!/usr/bin/env python3
"""
生成：富通光子_UFBG产品手册_中文版2.pptx
基于 product images 文件夹自动排版。
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

BASE = Path('/home/han/.openclaw/workspace/FBG Brochure')
IMG_DIR = BASE / 'product images'
OUT = BASE / '富通光子_UFBG产品手册_中文版2.pptx'

TITLE = '富通光子 UFBG 产品手册'
SUBTITLE = '中文版 2.0（自动生成）'


def add_title_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.5), Inches(1.0))
    tf = title_box.text_frame
    tf.text = TITLE
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT

    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(11.0), Inches(0.8))
    tf2 = sub_box.text_frame
    tf2.text = SUBTITLE
    p2 = tf2.paragraphs[0]
    p2.font.size = Pt(20)
    p2.alignment = PP_ALIGN.LEFT


def add_image_slide(prs: Presentation, img_path: Path, idx: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12.0), Inches(0.6))
    tf = title_box.text_frame
    tf.text = f'产品图片 {idx:02d}'
    p = tf.paragraphs[0]
    p.font.size = Pt(22)
    p.font.bold = True

    # 图片区域（16:9页面）
    left = Inches(0.7)
    top = Inches(1.0)
    max_w = Inches(11.9)
    max_h = Inches(5.5)

    # 先按宽放，再按高限制
    pic = slide.shapes.add_picture(str(img_path), left, top, width=max_w)
    if pic.height > max_h:
        pic.height = max_h
        # 按高度缩放后重新水平居中
        pic.left = int((prs.slide_width - pic.width) / 2)
    else:
        pic.left = int((prs.slide_width - max_w) / 2)

    # 底部文件名
    cap = slide.shapes.add_textbox(Inches(0.7), Inches(6.7), Inches(11.6), Inches(0.4))
    tfc = cap.text_frame
    tfc.text = img_path.name
    pc = tfc.paragraphs[0]
    pc.font.size = Pt(11)
    pc.alignment = PP_ALIGN.RIGHT


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)

    images = sorted([p for p in IMG_DIR.glob('*.jpg')])
    for i, img in enumerate(images, 1):
        add_image_slide(prs, img, i)

    prs.save(str(OUT))
    print(OUT)
    print(f'images_used={len(images)}')


if __name__ == '__main__':
    main()
