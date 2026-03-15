#!/usr/bin/env python3
"""
在“富通光子_UFBG产品手册_中文版.pptx”基础上生成“中文版2”：
- 保留原PPT全部页面和格式不变
- 仅追加“product images”中的图片
- 图片顺序与PDF采集顺序一致（按文件名 pageXX_imgYY）
- 在每页底部写入来源信息（页码/序号/文件名）
"""
from __future__ import annotations
import re
from pathlib import Path
from PIL import Image, ImageStat
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

BASE = Path('/home/han/.openclaw/workspace/FBG Brochure/dropbox_assets')
SRC_PPTX = BASE / '富通光子_UFBG产品手册_中文版.pptx'
IMG_DIR = BASE / 'product images'
OUT_PPTX = Path('/home/han/.openclaw/workspace/FBG Brochure/富通光子_UFBG产品手册_中文版2.pptx')

IMG_RE = re.compile(r'^page(\d+)_img(\d+)_', re.IGNORECASE)


def parse_order(name: str):
    m = IMG_RE.match(name)
    if m:
        return (int(m.group(1)), int(m.group(2)), name)
    return (10**9, 10**9, name)


def is_likely_logo(path: Path) -> bool:
    """简单logo过滤：极小图、长条横幅、近单色。"""
    try:
        im = Image.open(path).convert('RGB')
        w, h = im.size
        area = w * h
        if area < 60000:
            return True
        aspect = max(w / max(h, 1), h / max(w, 1))
        if aspect > 6:
            return True

        # 颜色变化很小 -> 可能是logo/图标
        stat = ImageStat.Stat(im)
        var = sum(stat.var) / 3
        if var < 120:
            return True
        return False
    except Exception:
        return False


def pick_images():
    exts = {'.jpg', '.jpeg', '.png', '.webp'}
    imgs = [p for p in IMG_DIR.iterdir() if p.is_file() and p.suffix.lower() in exts and p.name.lower().startswith('page')]
    imgs.sort(key=lambda p: parse_order(p.name))

    filtered = []
    for p in imgs:
        if not is_likely_logo(p):
            filtered.append(p)
    return filtered


def add_image_slide(prs: Presentation, img_path: Path, idx: int):
    # 使用空白版式，继承原模板主题
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    sw, sh = prs.slide_width, prs.slide_height

    # 顶部标题
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.5))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = f'产品图片补充 {idx:02d}'
    run.font.size = Pt(18)
    run.font.bold = True

    # 图片放置区
    left = Inches(0.6)
    top = Inches(0.85)
    max_w = Inches(12.1)
    max_h = Inches(5.9)

    pic = slide.shapes.add_picture(str(img_path), left, top, width=max_w)
    if pic.height > max_h:
        pic.height = max_h
    pic.left = int((sw - pic.width) / 2)
    pic.top = int((sh - pic.height) / 2)

    # 来源信息（根据文件名）
    m = IMG_RE.match(img_path.name)
    src = img_path.name
    if m:
        src = f'来源PDF: 第{int(m.group(1))}页 / 图{int(m.group(2))} | {img_path.name}'

    cap = slide.shapes.add_textbox(Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.4))
    ctf = cap.text_frame
    ctf.clear()
    cp = ctf.paragraphs[0]
    cp.alignment = PP_ALIGN.RIGHT
    cr = cp.add_run()
    cr.text = src
    cr.font.size = Pt(10)


def main():
    if not SRC_PPTX.exists():
        raise FileNotFoundError(f'未找到源PPT: {SRC_PPTX}')
    if not IMG_DIR.exists():
        raise FileNotFoundError(f'未找到图片目录: {IMG_DIR}')

    prs = Presentation(str(SRC_PPTX))
    images = pick_images()

    for i, img in enumerate(images, 1):
        add_image_slide(prs, img, i)

    OUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PPTX))
    print(f'out={OUT_PPTX}')
    print(f'images_added={len(images)}')
    print(f'slides_total={len(prs.slides)}')


if __name__ == '__main__':
    main()
