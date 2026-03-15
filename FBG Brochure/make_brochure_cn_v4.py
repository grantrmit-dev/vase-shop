#!/usr/bin/env python3
"""
在原始“富通光子_UFBG产品手册_中文版.pptx”基础上：
- 将 product images 中图片按 pageXX_imgYY 顺序生成图片页
- 与原文字页交错穿插（而不是全部追加到末尾）
- 输出：富通光子_UFBG产品手册_中文版2.pptx
"""
from __future__ import annotations
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

BASE = Path('/home/han/.openclaw/workspace/FBG Brochure/dropbox_assets')
SRC_PPTX = BASE / '富通光子_UFBG产品手册_中文版.pptx'
IMG_DIR = BASE / 'product images'
OUT_PPTX = Path('/home/han/.openclaw/workspace/FBG Brochure/富通光子_UFBG产品手册_中文版2.pptx')

IMG_RE = re.compile(r'^page(\d+)_img(\d+)_', re.IGNORECASE)


def parse_order(name: str):
    m = IMG_RE.match(name)
    if m:
        return (int(m.group(1)), int(m.group(2)), name)
    return (10**9, 10**9, name)


def pick_images():
    exts = {'.jpg', '.jpeg', '.png', '.webp'}
    imgs = [p for p in IMG_DIR.iterdir() if p.is_file() and p.suffix.lower() in exts and p.name.lower().startswith('page')]
    imgs.sort(key=lambda p: parse_order(p.name))
    return imgs


def add_image_slide(prs: Presentation, img_path: Path, idx: int):
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    sw, sh = prs.slide_width, prs.slide_height

    # 轻量标题（尽量不改变整体风格）
    tb = slide.shapes.add_textbox(Inches(0.45), Inches(0.18), Inches(12.2), Inches(0.45))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = f'产品图示 {idx:02d}'
    run.font.size = Pt(16)
    run.font.bold = True

    left = Inches(0.55)
    top = Inches(0.75)
    max_w = Inches(12.2)
    max_h = Inches(5.95)

    pic = slide.shapes.add_picture(str(img_path), left, top, width=max_w)
    if pic.height > max_h:
        pic.height = max_h
    pic.left = int((sw - pic.width) / 2)
    pic.top = int((sh - pic.height) / 2)

    # 底部来源
    m = IMG_RE.match(img_path.name)
    src = img_path.name
    if m:
        src = f'来源PDF: 第{int(m.group(1))}页 / 图{int(m.group(2))}'

    cap = slide.shapes.add_textbox(Inches(0.45), Inches(6.9), Inches(12.2), Inches(0.35))
    ctf = cap.text_frame
    ctf.clear()
    cp = ctf.paragraphs[0]
    cp.alignment = PP_ALIGN.RIGHT
    cr = cp.add_run()
    cr.text = src
    cr.font.size = Pt(10)


def reorder_interleaved(prs: Presentation, text_count: int, image_count: int):
    """将 slideId 列表重排为：文字页 + 若干图片页 + 文字页 ..."""
    sldIdLst = prs.slides._sldIdLst  # lxml element list
    ids = list(sldIdLst)

    text_ids = ids[:text_count]
    img_ids = ids[text_count:text_count + image_count]

    order = []
    img_i = 0
    remaining_imgs = image_count
    remaining_texts = text_count

    for t in text_ids:
        order.append(t)
        remaining_texts -= 1

        if remaining_imgs <= 0:
            continue

        # 均匀分配：每个文字页后插入 ceil(剩余图片/剩余文字页+1)
        slots = remaining_texts + 1
        k = (remaining_imgs + slots - 1) // slots
        for _ in range(k):
            if img_i < image_count:
                order.append(img_ids[img_i])
                img_i += 1
                remaining_imgs -= 1

    # 保险：若还有剩余图片，追加末尾
    while img_i < image_count:
        order.append(img_ids[img_i])
        img_i += 1

    # 应用新顺序
    for node in list(sldIdLst):
        sldIdLst.remove(node)
    for node in order:
        sldIdLst.append(node)


def main():
    if not SRC_PPTX.exists():
        raise FileNotFoundError(f'未找到源PPT: {SRC_PPTX}')

    images = pick_images()
    prs = Presentation(str(SRC_PPTX))
    text_count = len(prs.slides)

    for i, img in enumerate(images, 1):
        add_image_slide(prs, img, i)

    reorder_interleaved(prs, text_count, len(images))

    OUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PPTX))

    print(f'out={OUT_PPTX}')
    print(f'text_slides={text_count}')
    print(f'images_added={len(images)}')
    print(f'slides_total={len(prs.slides)}')


if __name__ == '__main__':
    main()
