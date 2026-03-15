#!/usr/bin/env python3
"""
基于原版“富通光子_UFBG产品手册_中文版.pptx”，在对应产品描述页内（文字旁）插入 product images 图片。
- 不改变原有版式/文字/表格
- 仅在对应页新增图片与细边框
- 输出：富通光子_UFBG产品手册_中文版2.pptx
"""
from __future__ import annotations
from pathlib import Path
import re
from collections import defaultdict
from PIL import Image
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

BASE = Path('/home/han/.openclaw/workspace/FBG Brochure/dropbox_assets')
SRC_PPTX = BASE / '富通光子_UFBG产品手册_中文版.pptx'
IMG_DIR = BASE / 'product images'
OUT_PPTX = Path('/home/han/.openclaw/workspace/FBG Brochure/富通光子_UFBG产品手册_中文版2.pptx')

IMG_RE = re.compile(r'^page(\d+)_img(\d+)_', re.IGNORECASE)


def parse_info(p: Path):
    m = IMG_RE.match(p.name)
    if not m:
        return None
    page = int(m.group(1))
    idx = int(m.group(2))
    try:
        with Image.open(p) as im:
            w, h = im.size
    except Exception:
        w, h = 0, 0
    return {
        'path': p,
        'page': page,
        'idx': idx,
        'w': w,
        'h': h,
        'area': w * h,
    }


def load_images():
    exts = {'.jpg', '.jpeg', '.png', '.webp'}
    items = []
    for p in IMG_DIR.iterdir():
        if p.is_file() and p.suffix.lower() in exts:
            info = parse_info(p)
            if info:
                items.append(info)
    return items


def pick_best(items, pages):
    cand = [x for x in items if x['page'] in pages]
    if not cand:
        return None
    # 选分辨率较高且非极端小图
    cand.sort(key=lambda x: (x['area'], x['w'], x['h']), reverse=True)
    return cand[0]['path']


def add_image_on_slide(slide, img_path: Path, x_cm=13.8, y_cm=2.0, w_cm=6.7, h_cm=4.3):
    # 先按宽放入，再按高限缩
    pic = slide.shapes.add_picture(str(img_path), Cm(x_cm), Cm(y_cm), width=Cm(w_cm))
    if pic.height > Cm(h_cm):
        # 按高度重新放一张（删除旧图）
        sp = pic._element
        sp.getparent().remove(sp)
        pic = slide.shapes.add_picture(str(img_path), Cm(x_cm), Cm(y_cm), height=Cm(h_cm))

    # 边框
    border = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        pic.left - Cm(0.05),
        pic.top - Cm(0.05),
        pic.width + Cm(0.10),
        pic.height + Cm(0.10),
    )
    border.fill.background()
    border.line.color.rgb = RGBColor(0x00, 0xB4, 0xD8)
    border.line.width = Pt(1.0)


def main():
    if not SRC_PPTX.exists():
        raise FileNotFoundError(f'源PPT不存在: {SRC_PPTX}')
    if not IMG_DIR.exists():
        raise FileNotFoundError(f'图片目录不存在: {IMG_DIR}')

    items = load_images()
    prs = Presentation(str(SRC_PPTX))

    # 按“产品描述页”映射 PDF 页段（可调）
    mapping = {
        5: range(4, 6),    # UFBG_S 示例1
        6: range(5, 7),    # UFBG_S 示例2
        9: range(7, 9),    # UFBG_W
        11: range(8, 10),  # Specialised
        13: range(10, 12), # 温度传感器
        14: range(12, 14), # 压力传感器
        15: range(14, 16), # 应变传感器
        16: range(16, 22), # 土压力传感器
        17: range(25, 33), # 解调仪
    }

    inserted = []
    for slide_idx, pages in mapping.items():
        if slide_idx < 1 or slide_idx > len(prs.slides):
            continue
        img = pick_best(items, set(pages))
        if not img:
            continue
        slide = prs.slides[slide_idx - 1]
        # 传感器页内容更满，右上图稍小
        if slide_idx in (13, 14, 15, 16, 17):
            add_image_on_slide(slide, img, x_cm=13.9, y_cm=1.9, w_cm=6.5, h_cm=3.8)
        else:
            add_image_on_slide(slide, img, x_cm=13.6, y_cm=2.1, w_cm=6.9, h_cm=4.4)
        inserted.append((slide_idx, img.name))

    OUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PPTX))

    print(f'out={OUT_PPTX}')
    print(f'inserted_count={len(inserted)}')
    for s, n in inserted:
        print(f'slide{s}: {n}')


if __name__ == '__main__':
    main()
